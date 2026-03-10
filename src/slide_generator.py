"""
slide_generator.py  — FIXED
-----------------------------
Fixes:
1. Function name is generate_all_slides (consistent everywhere)
2. Uses plan_path parameter name (consistent with other modules)
3. Handles merged subtopics — shows all merged titles on title slide
"""

import json
import re
import subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Theme colours ─────────────────────────────────────────────
DARK  = RGBColor(0x0a, 0x0d, 0x14)
BLUE  = RGBColor(0x00, 0xd4, 0xff)
WHITE = RGBColor(0xe0, 0xe6, 0xf0)
GREY  = RGBColor(0x4a, 0x55, 0x70)
GREEN = RGBColor(0x00, 0xff, 0x9d)
DARK2 = RGBColor(0x06, 0x0f, 0x1a)


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _txt(slide, text, l, t, w, h, size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def make_pptx(title: str, topic: str, content: str, merged_from: list, order: int, out: str):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Title slide ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, DARK)

    # Topic label
    _txt(s, topic.upper(), Inches(0.5), Inches(0.3), Inches(12), Inches(0.4),
         size=11, color=GREY)

    # Accent bar
    bar = s.shapes.add_shape(1, Inches(0.5), Inches(0.9), Inches(0.06), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

    # Main title
    _txt(s, title, Inches(0.8), Inches(0.85), Inches(11), Inches(1.4),
         size=36, bold=True)

    # Merged subtopics hint
    if len(merged_from) > 1:
        sub = "Covers: " + "  ·  ".join(merged_from)
        _txt(s, sub, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.5),
             size=13, color=GREY)

    # ── Content slides (5 bullets each) ──────────────────────
    sents  = [x.strip() for x in re.split(r"(?<=[.!?])\s+", content) if len(x.strip()) > 15]
    chunks = [sents[i:i+5] for i in range(0, max(1, len(sents)), 5)]

    for chunk in chunks:
        s = prs.slides.add_slide(blank)
        _bg(s, DARK)
        # top bar
        bar2 = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
        bar2.fill.solid(); bar2.fill.fore_color.rgb = BLUE; bar2.line.fill.background()
        # title
        _txt(s, title, Inches(0.5), Inches(0.12), Inches(12), Inches(0.45), size=12, color=GREY)

        top = Inches(0.8)
        for line in chunk:
            tb = s.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.9))
            tb.text_frame.word_wrap = True
            p  = tb.text_frame.paragraphs[0]
            r1 = p.add_run();  r1.text = "→  ";  r1.font.size = Pt(18); r1.font.color.rgb = BLUE
            r2 = p.add_run();  r2.text = line[:150]; r2.font.size = Pt(18); r2.font.color.rgb = WHITE
            top += Inches(1.05)

    # ── Summary slide ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, DARK2)
    _txt(s, "KEY TAKEAWAY", Inches(0.5), Inches(0.3), Inches(12), Inches(0.45),
         size=12, bold=True, color=BLUE)
    box = s.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(12.3), Inches(3.0))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x06, 0x1e, 0x14)
    box.line.color.rgb = GREEN; box.line.width = Pt(1.5)
    summary = sents[-1] if sents else content[:200]
    _txt(s, summary[:300], Inches(0.9), Inches(1.4), Inches(11.5), Inches(2.4), size=20)

    prs.save(out)


def pptx_to_pngs(pptx_path: str, out_dir: str) -> list:
    Path(out_dir).mkdir(parents=True, exist_ok=True)
    subprocess.run([
        "libreoffice", "--headless", "--convert-to", "png",
        "--outdir", out_dir, pptx_path
    ], check=True, capture_output=True)
    return sorted(str(p) for p in Path(out_dir).glob("*.png"))


def generate_all_slides(
    plan_path:  str = "output/json/lecture_plan.json",
    slides_dir: str = "output/slides"
):
    if not Path(plan_path).exists():
        raise FileNotFoundError(f"Not found: {plan_path}")

    data      = json.loads(Path(plan_path).read_text(encoding="utf-8"))
    subtopics = data.get("subtopics", [])

    print(f"\n📊 Building slides for {len(subtopics)} groups…")

    for i, st in enumerate(subtopics):
        title       = st.get("title", f"group_{i}")
        safe        = re.sub(r"[^\w\-]", "_", title)[:50]
        sub_dir     = Path(slides_dir) / f"{i:02d}_{safe}"
        sub_dir.mkdir(parents=True, exist_ok=True)
        pptx_p      = str(sub_dir / "slides.pptx")
        png_dir     = str(sub_dir / "pngs")
        merged_from = st.get("merged_from", [title])

        print(f"   [{i+1}/{len(subtopics)}] {title}")

        make_pptx(
            title       = title,
            topic       = st.get("topic", ""),
            content     = st.get("content", ""),
            merged_from = merged_from,
            order       = i,
            out         = pptx_p
        )

        pngs = pptx_to_pngs(pptx_p, png_dir)
        st["slides_dir"]  = png_dir
        st["slide_count"] = len(pngs)
        print(f"      → {len(pngs)} slides")

    data["subtopics"] = subtopics
    Path(plan_path).write_text(
        json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    print(f"\n✅ Slides done → {slides_dir}")


if __name__ == "__main__":
    generate_all_slides()